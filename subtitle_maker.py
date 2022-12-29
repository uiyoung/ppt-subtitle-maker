from pptx import Presentation


def generate_ppt(songs):
    # open template
    presentation = Presentation('template.pptx')

    for song in songs:
        song_title = song[1]
        lyrics = song[2].split('\n\n')
        lyrics = list(map(lambda s: s.strip(), lyrics))  # 각줄의 양끝 개행문자 제거

        title_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_layout)

        title_placeholder = slide.shapes.title
        title_placeholder.text = song_title

        subtitle_layout = presentation.slide_layouts[1]

        # genrate subtitle slides with lyrics
        for subtitle in lyrics:
            slide = presentation.slides.add_slide(subtitle_layout)
            title_placeholder = slide.shapes.title
            title_placeholder.text = subtitle

    # save
    presentation.save('./ppt/' + song_title + '.pptx')

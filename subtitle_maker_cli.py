from pptx import Presentation
import sys

# get arguments
argument = sys.argv


# read lyrics text file
f = open('./txt/' + argument[1] + '.txt', 'r', encoding='UTF-8')
lines = f.readlines()
lines = list(map(lambda s: s.strip(), lines))  # 각줄의 양끝 개행문자 제거
# lines = list(filter(lambda s: s not in ['', '\n', '\u200b\n' ',\u200b'], lines))  # remove empty line

# insert '' to every 3rd line
for i in range(2, len(lines), 3):
    if lines[i] != '':
        lines.insert(i, '')


# combine 3 elements from lines
lyrics = []
for i in range(0, len(lines), 3):
    lyrics.append('\n'.join(lines[i:i+3]).strip())


# open template
presentation = Presentation('template.pptx')

# set title
song_title = argument[1]

title_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_layout)

title_placeholder = slide.shapes.title
title_placeholder.text = song_title

subtitle_layout = presentation.slide_layouts[1]

# genrate subtitle slides with lyrics
for subtitle in lyrics:
    slide = presentation.slides.add_slide(subtitle_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = subtitle


# save
presentation.save('./ppt/' + song_title + '.pptx')

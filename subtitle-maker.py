from pptx import Presentation
import sys

# get arguments
argument = sys.argv


# read lyrics text file
f = open(argument[1] + '.txt', 'r', encoding='UTF-8')
lines = f.readlines()
lines = list(filter(lambda s: s != ('\u200b\n' and '\n'), lines))  # remove emptyline

# combine 2 elements from lines
lyrics = []
for i in range(0, len(lines), 2):
    lyrics.append(''.join(lines[i:i+2]))

# 개행문자 제거
lyrics = list(map(lambda s: s.strip(), lyrics))


# open template
presentation = Presentation('template.pptx')

# set title
song_title = argument[1]
title_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_layout)

title_placeholder = slide.shapes.title
title_placeholder.text = song_title

subtitle_layout = presentation.slide_layouts[1]

# make subtitle slides
for i in range(len(lyrics)):
    slide = presentation.slides.add_slide(subtitle_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = lyrics[i]


# save
presentation.save(song_title + '.pptx')
